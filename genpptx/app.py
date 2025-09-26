# app.py
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import os
import re
import json
from dotenv import load_dotenv

load_dotenv()
app = Flask(__name__)

# ===== Azure OpenAI 設定 =====
AI_AVAILABLE = False
client = None
try:
    from openai import AzureOpenAI
    api_key = os.getenv("AZURE_OPENAI_API_KEY")
    endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
    api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-01")
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o-mini")  # デプロイ名を環境変数で

    if api_key and endpoint:
        client = AzureOpenAI(
            api_key=api_key,
            azure_endpoint=endpoint,
            api_version=api_version
        )
        AI_AVAILABLE = True
        print("[INFO] Azure OpenAI enabled.")
    else:
        print("[WARN] Azure OpenAI credentials not found. Fallback mode.")
except Exception as e:
    print(f"[WARN] Azure OpenAI init error: {e}. Fallback mode.")

# ===== ユーティリティ =====
PRIMARY_RGB = RGBColor(32, 89, 167)     # 深めのブルー（見出し・アクセント）
ACCENT_RGB = RGBColor(237, 242, 248)    # 薄いグレー/ブルー（背景ブロック）
TEXT_RGB   = RGBColor(25, 25, 25)
SUBTEXT_RGB= RGBColor(90, 98, 110)

JP_FONTS = ["Yu Gothic UI", "Yu Gothic", "Meiryo", "MS PGothic", "Segoe UI"]

def set_font(run, size_pt=24, bold=False, color=TEXT_RGB):
    if run is None:
        return
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    # 日本語フォント優先
    for f in JP_FONTS:
        run.font.name = f
        break

def normalize_text(s: str) -> str:
    s = s.replace('\u3000', ' ').strip()  # 全角スペース→半角
    s = re.sub(r'[ \t]+', ' ', s)
    return s

def shorten_bullet(text: str, max_chars: int = 45) -> str:
    t = normalize_text(text)
    if len(t) <= max_chars:
        return t
    # 句点や読点、ダッシュで上手に省略
    for token in ['。', '、', ' - ', ' – ', ' — ', ';', '：', ':']:
        idx = t.find(token)
        if 0 < idx <= max_chars:
            return t[:idx].strip() + '…'
    return t[:max_chars].rstrip() + '…'

def safe_ascii_filename(name: str, default="meeting"):
    s = re.sub(r'[^\w\-.]+', '_', name, flags=re.UNICODE)
    return s or default

def strip_code_fence(s: str) -> str:
    return re.sub(r"^```(?:json)?|```$", "", s.strip(), flags=re.MULTILINE)

# ====== AI 解析 ======
def parse_meeting_minutes(minutes_text: str) -> dict:
    """
    議事録テキスト→構造化データ（スライド設計用）
    """
    minutes_text = minutes_text or ""
    if AI_AVAILABLE and client:
        try:
            system = (
                "あなたは上級のB2Bプリセールスです。"
                "入力の議事録を、ビジネスプレゼンの分かりやすい流れに適したJSONへ構造化してください。"
                "箇条書きは短く簡潔に（1行・名詞止めや体言止めを優先）。"
                "可能なら冗長表現を圧縮してください。"
            )
            user_prompt = f"""
以下の議事録から、スライド用の構造化JSONを返してください。
スキーマ:
{{
  "company_name": "string",
  "meeting_date": "string",
  "title": "string",
  "agenda": ["string"],
  "sections": [
    {{
      "title": "string",
      "bullets": ["string"],
      "notes": ["string"]
    }}
  ],
  "challenges": ["string"],
  "needs": ["string"],
  "next_actions": ["string"],
  "bant": {{"budget":"string", "authority":"string", "need":"string", "timeline":"string"}},
  "summary": ["string"]
}}
要件:
- 箇条書きは45文字程度で簡潔に。
- 無い要素は空配列/空文字でOK。
- 日本語で返す。
議事録:
{minutes_text}
"""
            resp = client.chat.completions.create(
                model=os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o-mini"),
                temperature=0.2,
                response_format={"type": "json_object"},
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": user_prompt},
                ],
            )
            content = resp.choices[0].message.content
            content = strip_code_fence(content)
            data = json.loads(content)
            return data
        except Exception as e:
            print(f"[WARN] AI parse error: {e}. Use fallback.")
            return parse_meeting_minutes_fallback(minutes_text)
    else:
        print("[INFO] AI disabled. Use fallback.")
        return parse_meeting_minutes_fallback(minutes_text)

def parse_meeting_minutes_fallback(minutes_text: str) -> dict:
    """
    AI無し時：Markdown風見出し（# / ## / ###）に基づくセクション抽出
    """
    lines = [normalize_text(l) for l in (minutes_text or "").splitlines()]
    company_name = "お客様"
    meeting_date = ""
    title = ""
    sections = []
    current = None

    for line in lines:
        if not line:
            continue
        # 会社名 or タイトル
        if ("株式会社" in line or "会社" in line) and line.startswith("#") and not line.startswith("##"):
            company_name = line.lstrip("# ").strip()
            if not title:
                title = f"{company_name} 打ち合わせ要約"
            continue
        if re.search(r"(日時|実施日)[：:]\s*\d{4}年\d{1,2}月\d{1,2}日", line):
            m = re.search(r"(\d{4}年\d{1,2}月\d{1,2}日)", line)
            meeting_date = m.group(1) if m else meeting_date
            continue
        # セクション
        if line.startswith("##"):
            if current:
                sections.append(current)
            current = {"title": line.lstrip("# ").strip(), "bullets": [], "notes": []}
            continue
        if line.startswith("###"):
            if current:
                st = line.lstrip("# ").strip()
                current["bullets"].append(f"【{shorten_bullet(st)}】")
            continue
        if re.match(r"^[-*•]\s+", line):
            if current:
                current["bullets"].append(shorten_bullet(re.sub(r"^[-*•]\s+", "", line)))
            continue
        # その他はノーツへ
        if current:
            current["notes"].append(shorten_bullet(line, max_chars=60))
    if current:
        sections.append(current)

    if not title:
        title = f"{company_name} 打ち合わせ要約"
    if not meeting_date:
        meeting_date = "2025年7月9日"  # デフォルト

    return {
        "company_name": company_name,
        "meeting_date": meeting_date,
        "title": title,
        "agenda": [s["title"].replace("【", "").replace("】", "") for s in sections][:6],
        "sections": sections,
        "challenges": [],
        "needs": [],
        "next_actions": [],
        "bant": {"budget": "", "authority": "", "need": "", "timeline": ""},
        "summary": [],
    }

# ====== PPT 生成 ======
def create_meeting_summary_ppt(parsed: dict) -> io.BytesIO:
    """
    構造化データ→PPTX生成
    """
    template_path = os.path.join(os.path.dirname(__file__), "image", "tempppt.pptx")
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print("[INFO] Use template:", template_path)
    else:
        prs = Presentation()
        # デフォルトのスライドサイズや簡易テーマは pptx 既定を利用

    def get_layout(index_fallback=1):
        try:
            return prs.slide_layouts[index_fallback]
        except Exception:
            return prs.slide_layouts[0]

    def add_title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(get_layout(0))
        if slide.shapes.title:
            slide.shapes.title.text = title
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 40, True, PRIMARY_RGB)
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        # 副題 - 安全なプレースホルダーアクセス
        try:
            if len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                ph.text = subtitle
                p = ph.text_frame.paragraphs[0]
                if p.runs:
                    set_font(p.runs[0], 20, False, SUBTEXT_RGB)
        except (KeyError, IndexError):
            # プレースホルダーが存在しない場合はテキストボックスを追加
            if subtitle:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
                textbox.text = subtitle
                p = textbox.text_frame.paragraphs[0]
                if p.runs:
                    set_font(p.runs[0], 20, False, SUBTEXT_RGB)

    def add_section_divider(title: str):
        # セクション仕切り
        slide = prs.slides.add_slide(get_layout(5 if len(prs.slide_layouts) > 5 else 1))
        if slide.shapes.title:
            slide.shapes.title.text = title
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 36, True, PRIMARY_RGB)

    def ensure_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    def add_bullets_block(shape, items, font_size=20, bullet_char="•", top_padding=False):
        tf = shape.text_frame
        tf.clear()
        if not items:
            return
        for i, it in enumerate(items):
            if i == 0:
                tf.text = f"{bullet_char} {it}"
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = f"{bullet_char} {it}"
            p.level = 0
            if p.runs:
                # フォントサイズを内容の長さで少し調整
                sz = font_size - 2 if len(it) > 40 else font_size
                set_font(p.runs[0], sz, False, TEXT_RGB)
        # 行間
        for p in tf.paragraphs:
            p.line_spacing = 1.15

    def add_content_slide(title: str, bullets: list, notes: list):
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = title
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 30, True, PRIMARY_RGB)

        # 箇条書きが多ければ2カラム
        items = bullets or []
        max_per_slide = 10
        if len(items) > max_per_slide:
            # 過剰分は次スライドへ
            chunk = items[:max_per_slide]
            rest = items[max_per_slide:]
            add_bullets_block(ensure_textbox(slide, 1.0, 2.0, 8.0, 4.8), chunk)
            if notes:
                add_bullets_block(ensure_textbox(slide, 1.0, 6.9, 8.0, 1.5), [f"補足: {n}" for n in notes[:2]], 16, "―")
            # 残りを再帰的に別スライド
            add_content_slide(f"{title}（続き）", rest, notes[2:] if notes else [])
            return

        # 2カラム条件（>8個）
        if len(items) >= 8:
            half = (len(items) + 1) // 2
            left_items = items[:half]
            right_items = items[half:]
            add_bullets_block(ensure_textbox(slide, 0.8, 2.0, 4.2, 5.2), left_items)
            add_bullets_block(ensure_textbox(slide, 5.2, 2.0, 4.2, 5.2), right_items)
        else:
            add_bullets_block(ensure_textbox(slide, 1.0, 2.0, 8.0, 5.5), items)

        # 補足ノートを下段に控えめ表示
        if notes:
            box = ensure_textbox(slide, 1.0, 7.1, 8.0, 1.2)
            tf = box.text_frame
            tf.clear()
            tf.text = "補足"
            p0 = tf.paragraphs[0]
            if p0.runs:
                set_font(p0.runs[0], 14, True, SUBTEXT_RGB)
            for n in notes[:3]:
                p = tf.add_paragraph()
                p.text = f"– {n}"
                if p.runs: set_font(p.runs[0], 12, False, SUBTEXT_RGB)

    def add_agenda_slide(items: list):
        if not items:
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "アジェンダ"
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 30, True, PRIMARY_RGB)
        add_bullets_block(ensure_textbox(slide, 1.0, 2.0, 8.0, 5.5), [shorten_bullet(i) for i in items], 22, "■")

    def add_bant_slide(bant: dict):
        if not bant:
            return
        if not any([bant.get(k) for k in ("budget", "authority", "need", "timeline")]):
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "BANT（商談情報）"
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 30, True, PRIMARY_RGB)
        # テーブル風レイアウト
        box = ensure_textbox(slide, 1.0, 2.0, 8.0, 5.5)
        tf = box.text_frame
        tf.clear()

        rows = [
            ("Budget（予算）", bant.get("budget", "")),
            ("Authority（決裁）", bant.get("authority", "")),
            ("Need（ニーズ）", bant.get("need", "")),
            ("Timeline（時期）", bant.get("timeline", "")),
        ]
        for i, (k, v) in enumerate(rows):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = k
            if p.runs: set_font(p.runs[0], 18, True, PRIMARY_RGB)
            p.line_spacing = 1.1
            p2 = tf.add_paragraph()
            p2.text = f"{v or '—'}"
            if p2.runs: set_font(p2.runs[0], 18, False, TEXT_RGB)
            p2.space_after = Pt(10)

    def add_next_actions_slide(items: list):
        if not items:
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "ネクストアクション"
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 30, True, PRIMARY_RGB)
        add_bullets_block(ensure_textbox(slide, 1.0, 2.0, 8.0, 5.0), [shorten_bullet(i) for i in items], 22, "□")

    def add_summary_slide(items: list):
        if not items:
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "まとめ"
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 30, True, PRIMARY_RGB)
        add_bullets_block(ensure_textbox(slide, 1.0, 2.0, 8.0, 5.5), [shorten_bullet(i) for i in items], 22, "●")

    # ===== 実装：スライド生成フロー =====
    title = parsed.get("title") or "打ち合わせ要約"
    company = parsed.get("company_name") or ""
    date = parsed.get("meeting_date") or ""
    subtitle = " / ".join([s for s in [company, f"{date} 実施" if date else ""] if s])

    add_title_slide(title, subtitle)
    add_agenda_slide(parsed.get("agenda", []))

    # セクションごと
    sections = parsed.get("sections", [])
    for s in sections:
        stitle = s.get("title") or ""
        bullets = [shorten_bullet(b) for b in (s.get("bullets") or [])]
        notes = s.get("notes") or []
        # セクション見出しスライド（大見出しに見える場合のみ）
        if stitle and (len(bullets) == 0 and len(notes) == 0):
            add_section_divider(stitle)
        # 通常コンテンツ
        add_content_slide(stitle or "セクション", bullets, notes)

    # BANT／課題・ニーズ／ネクストアクション／まとめ
    add_bant_slide(parsed.get("bant", {}))

    # 課題・ニーズがあれば1枚に統合
    challenges = [shorten_bullet(x) for x in parsed.get("challenges", [])]
    needs = [shorten_bullet(x) for x in parsed.get("needs", [])]
    if challenges or needs:
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "課題とニーズ"
            set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 30, True, PRIMARY_RGB)
        add_bullets_block(ensure_textbox(slide, 0.8, 2.0, 4.2, 5.2), ["【課題】"] + challenges if challenges else [], 22, "•")
        add_bullets_block(ensure_textbox(slide, 5.2, 2.0, 4.2, 5.2), ["【ニーズ】"] + needs if needs else [], 22, "•")

    add_next_actions_slide(parsed.get("next_actions", []))
    add_summary_slide(parsed.get("summary", []))

    # 末尾に Thanks スライド
    slide = prs.slides.add_slide(get_layout(1))
    if slide.shapes.title:
        slide.shapes.title.text = "ご清聴ありがとうございました"
        set_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], 32, True, PRIMARY_RGB)

    # 出力
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ===== Flask ルーティング =====
@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    minutes_text = request.form.get("minutes_text", "")
    parsed = parse_meeting_minutes(minutes_text)
    ppt_file = create_meeting_summary_ppt(parsed)

    company_name = parsed.get("company_name", "meeting") or "meeting"
    fname = f"{safe_ascii_filename(company_name)}_summary.pptx"

    return send_file(
        ppt_file,
        as_attachment=True,
        download_name=fname,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5001))
    app.run(debug=False, host="0.0.0.0", port=port)
