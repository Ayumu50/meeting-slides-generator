from flask import Flask, render_template, request, send_file, redirect, url_for, session
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import os
import re
import json
import tempfile
import uuid
from datetime import datetime
from dotenv import load_dotenv

load_dotenv()
app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'your-secret-key-here')

# 生成されたファイルを一時保存するディレクトリ
TEMP_DIR = os.path.join(os.path.dirname(__file__), 'temp_files')
os.makedirs(TEMP_DIR, exist_ok=True)

# ===== Azure OpenAI 設定 =====
try:
    from openai import AzureOpenAI
    api_key = os.getenv("AZURE_OPENAI_API_KEY")
    endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
    api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-01")
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o-mini")  # デプロイ名を環境変数で

    if not api_key or not endpoint:
        raise ValueError("Azure OpenAI credentials are required. Please set AZURE_OPENAI_API_KEY and AZURE_OPENAI_ENDPOINT environment variables.")
    
    client = AzureOpenAI(
        api_key=api_key,
        azure_endpoint=endpoint,
        api_version=api_version
    )
    print("[INFO] Azure OpenAI enabled.")
except Exception as e:
    print(f"[ERROR] Azure OpenAI initialization failed: {e}")
    raise

# ===== ユーティリティ =====
# カラールール: background / text / main / accent
BACKGROUND_RGB = RGBColor(255, 255, 255)  # 背景は白を前提
PRIMARY_RGB = RGBColor(32, 89, 167)     # メインカラー（見出し・アクセント）
ACCENT_RGB = RGBColor(237, 242, 248)    # アクセント（淡い塗り）
TEXT_RGB   = RGBColor(25, 25, 25)
SUBTEXT_RGB= RGBColor(90, 98, 110)

JP_FONTS = ["Yu Gothic UI", "Yu Gothic", "Meiryo", "MS PGothic", "Segoe UI"]

# 設計ポリシー（okunote の 9つのコツを反映）
MAX_BULLETS_PER_SLIDE = 6     # 1スライドあたりの箇条は極力6件以下
MAX_CHARS_PER_BULLET = 45     # 箇条は短く
ONE_MESSAGE_POLICY = True     # ワンスライド・ワンメッセージ


def set_font(run, size_pt=24, bold=False, color=TEXT_RGB, align_left=True):
    if run is None:
        return
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    # 日本語フォント優先
    for f in JP_FONTS:
        run.font.name = f
        break


def normalize_text(s: str) -> str:
    s = s or ""
    s = s.replace('\u3000', ' ').strip()  # 全角スペース→半角
    s = re.sub(r'[ \t]+', ' ', s)
    return s


def shorten_bullet(text: str, max_chars: int = MAX_CHARS_PER_BULLET) -> str:
    t = normalize_text(text)
    if len(t) <= max_chars:
        return t
    # 句点や読点、ダッシュで上手に省略
    for token in ['。', '、', ' - ', ' – ', ' — ', ';', '：', ':']:
        idx = t.find(token)
        if 0 < idx <= max_chars:
            return t[:idx].strip() + '…'
    return t[:max_chars].rstrip() + '…'


def safe_ascii_filename(name: str, default="meeting"):
    s = re.sub(r'[^\w\-.]+', '_', name, flags=re.UNICODE)
    return s or default


def strip_code_fence(s: str) -> str:
    return re.sub(r"^```(?:json)?|```$", "", s.strip(), flags=re.MULTILINE)

# ====== AI 解析 ======
def parse_meeting_minutes(minutes_text: str) -> dict:
    """
    議事録テキスト→構造化データ（スライド設計用）
    生成AIのみを使用して解析を行います。生成には5分ほど時間がかかります。
    """
    minutes_text = minutes_text or ""
    
    if not minutes_text.strip():
        raise ValueError("議事録テキストが入力されていません。")
    
    try:
        system = (
            "あなたは上級のB2Bプリセールスです。"
            "入力の議事録を、ビジネスプレゼンの分かりやすい流れに適したJSONへ構造化してください。"
            "箇条書きは短く簡潔に（1行・名詞止めや体言止めを優先）。"
            "可能なら冗長表現を圧縮してください。"
        )
        user_prompt = f"""
以下の議事録から、スライド用の構造化JSONを返してください。
スキーマ:
{{
  "company_name": "string",
  "meeting_date": "string",
  "title": "string",
  "agenda": ["string"],
  "sections": [
    {{
      "title": "string",
      "bullets": ["string"],
      "notes": ["string"]
    }}
  ],
  "challenges": ["string"],
  "needs": ["string"],
  "next_actions": ["string"],
  "bant": {{"budget":"string", "authority":"string", "need":"string", "timeline":"string"}},
  "summary": ["string"]
}}
要件:
- 箇条書きは45文字程度で簡潔に。
- 無い要素は空配列/空文字でOK。
- 日本語で返す。
議事録:
{minutes_text}
"""
        resp = client.chat.completions.create(
            model=os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-4o-mini"),
            temperature=1.0,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": user_prompt},
            ],
        )
        content = resp.choices[0].message.content
        content = strip_code_fence(content)
        data = json.loads(content)
        return data
    except Exception as e:
        print(f"[ERROR] AI parse error: {e}")
        raise ValueError(f"議事録の解析に失敗しました: {str(e)}")




# ====== PPT 生成 ======
def create_meeting_summary_ppt(parsed: dict) -> io.BytesIO:
    """
    構造化データ→PPTX生成

    変更点（okunote の 9つのコツを反映）:
    - ワンスライド・ワンメッセージ原則の適用
    - 箇条は最大6件に制限、長い箇条は短縮して別スライド化
    - 視線を意識して重要な要素を上へ配置
    - 色の役割（背景・文字・メイン・アクセント）を厳格化
    - 左揃えの徹底、行間・余白の調整
    """
    template_path = os.path.join(os.path.dirname(__file__), "image", "tempppt.pptx")
    try:
        if os.path.exists(template_path):
            prs = Presentation(template_path)
            print("[INFO] Using template:", template_path)
        else:
            prs = Presentation()
            print("[WARNING] Template not found, using default presentation")
    except Exception as e:
        print(f"[ERROR] Failed to load template: {e}")
        prs = Presentation()
        print("[INFO] Using default presentation as fallback")

    def get_layout(index_fallback=1):
        try:
            return prs.slide_layouts[index_fallback]
        except Exception:
            return prs.slide_layouts[0]

    def add_title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(get_layout(0))
        # タイトルの背景に淡いアクセント矩形を置き、タイトルを上部に寄せる（視線を上へ）
        try:
            # タイトル背景
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)).fill.solid()
            rect = slide.shapes[-1]
            rect.fill.fore_color.rgb = ACCENT_RGB
            rect.line.fill.background()
            # 透明度はpptxで直接制御しにくいため薄い色を使う
        except Exception:
            pass

        if slide.shapes.title:
            slide.shapes.title.text = title
            t = slide.shapes.title.text_frame.paragraphs[0]
            if t.runs:
                set_font(t.runs[0], 40, True, PRIMARY_RGB)
            t.alignment = PP_ALIGN.LEFT

        # 副題
        try:
            if len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                ph.text = subtitle
                p = ph.text_frame.paragraphs[0]
                if p.runs:
                    set_font(p.runs[0], 18, False, SUBTEXT_RGB)
                p.alignment = PP_ALIGN.LEFT
        except (KeyError, IndexError):
            if subtitle:
                textbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9), Inches(0.6))
                textbox.text = subtitle
                p = textbox.text_frame.paragraphs[0]
                if p.runs:
                    set_font(p.runs[0], 18, False, SUBTEXT_RGB)
                p.alignment = PP_ALIGN.LEFT

    def add_section_divider(title: str):
        slide = prs.slides.add_slide(get_layout(5 if len(prs.slide_layouts) > 5 else 1))
        if slide.shapes.title:
            slide.shapes.title.text = title
            t = slide.shapes.title.text_frame.paragraphs[0]
            if t.runs:
                set_font(t.runs[0], 36, True, PRIMARY_RGB)
            t.alignment = PP_ALIGN.LEFT

    def ensure_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    def add_bullets_block(shape, items, font_size=20, bullet_char="•"):
        tf = shape.text_frame
        tf.clear()
        if not items:
            return
        # 最初の段落は見出し扱いにしない（箇条のみにする）
        for i, it in enumerate(items):
            if i == 0:
                tf.text = f"{bullet_char} {it}"
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = f"{bullet_char} {it}"
            p.level = 0
            # 左揃え
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                # フォントサイズを内容の長さで少し調整
                sz = font_size - 2 if len(it) > 40 else font_size
                set_font(p.runs[0], sz, False, TEXT_RGB)
            p.line_spacing = 1.2
            p.space_after = Pt(6)

    def promote_key_message(slide, message: str):
        """重要メッセージをスライド上部に大きく表示する（ワンスライド・ワンメッセージ）"""
        box = ensure_textbox(slide, 0.6, 1.6, 9.0, 1.0)
        tf = box.text_frame
        tf.clear()
        tf.text = message
        p = tf.paragraphs[0]
        if p.runs:
            set_font(p.runs[0], 22, True, PRIMARY_RGB)
        p.alignment = PP_ALIGN.LEFT

    def add_content_slide(title: str, bullets: list, notes: list):
        # ワンスライド・ワンメッセージの適用:
        # - bullets が 1 件で十分に要約されている場合はタイトル＋1行で完結
        # - bullets が複数だが各々が短くて関連性が高い場合は 1 スライドに最大 MAX_BULLETS_PER_SLIDE
        # - 長い箇条や1件ごとにメッセージ性が高い箇条は個別スライドへ分割
        items = bullets or []
        # 重要な要素をスライド上部へ
        if len(items) == 1:
            slide = prs.slides.add_slide(get_layout(1))
            if slide.shapes.title:
                slide.shapes.title.text = title
                t = slide.shapes.title.text_frame.paragraphs[0]
                if t.runs:
                    set_font(t.runs[0], 30, True, PRIMARY_RGB)
                t.alignment = PP_ALIGN.LEFT
            # 1つのメッセージを大きめに表示
            promote_key_message(slide, shorten_bullet(items[0], max_chars=80))
            # 補足ノート
            if notes:
                box = ensure_textbox(slide, 0.6, 3.0, 9.0, 2.5)
                add_bullets_block(box, [shorten_bullet(n, max_chars=80) for n in notes], 16)
            return

        # 分割対象（各箇条がそれぞれメッセージ）
        # 判定: 箇条が短いが個別に伝える価値がある（例：箇条ごとに句点含む or 長さ>max）
        separate_slides = []
        remaining = []
        for it in items:
            clean = normalize_text(it)
            # 個々の箇条が独立したメッセージと判断する条件
            if len(clean) > MAX_CHARS_PER_BULLET * 0.9 or '\n' in clean:
                separate_slides.append(clean)
            else:
                remaining.append(clean)

        # まず残りを1スライドにまとめる（ただし件数が多ければ分割）
        if remaining:
            # chunking
            chunks = [remaining[i:i+MAX_BULLETS_PER_SLIDE] for i in range(0, len(remaining), MAX_BULLETS_PER_SLIDE)]
            for idx, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(get_layout(1))
                slide_title = title if idx == 0 else f"{title}（続き）"
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                    t = slide.shapes.title.text_frame.paragraphs[0]
                    if t.runs:
                        set_font(t.runs[0], 30, True, PRIMARY_RGB)
                    t.alignment = PP_ALIGN.LEFT
                # 2カラムを検討（chunk 数が多ければ2カラム）
                if len(chunk) >= 5:
                    half = (len(chunk) + 1) // 2
                    add_bullets_block(ensure_textbox(slide, 0.6, 1.9, 4.4, 4.8), [shorten_bullet(i) for i in chunk[:half]])
                    add_bullets_block(ensure_textbox(slide, 5.2, 1.9, 4.4, 4.8), [shorten_bullet(i) for i in chunk[half:]])
                else:
                    add_bullets_block(ensure_textbox(slide, 0.6, 1.9, 9.0, 4.8), [shorten_bullet(i) for i in chunk])
                if notes and idx == 0:
                    box = ensure_textbox(slide, 0.6, 6.8, 9.0, 1.0)
                    add_bullets_block(box, [shorten_bullet(n, max_chars=80) for n in notes[:2]], 16)

        # separate_slides は1件ずつ強調
        for msg in separate_slides:
            slide = prs.slides.add_slide(get_layout(1))
            if slide.shapes.title:
                slide.shapes.title.text = title
                t = slide.shapes.title.text_frame.paragraphs[0]
                if t.runs:
                    set_font(t.runs[0], 30, True, PRIMARY_RGB)
                t.alignment = PP_ALIGN.LEFT
            promote_key_message(slide, shorten_bullet(msg, max_chars=120))

    def add_agenda_slide(items: list):
        if not items:
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "アジェンダ"
            t = slide.shapes.title.text_frame.paragraphs[0]
            if t.runs:
                set_font(t.runs[0], 30, True, PRIMARY_RGB)
            t.alignment = PP_ALIGN.LEFT
        add_bullets_block(ensure_textbox(slide, 0.6, 1.9, 9.0, 4.8), [shorten_bullet(i) for i in items], 20, "■")

    def add_bant_slide(bant: dict):
        if not bant:
            return
        if not any([bant.get(k) for k in ("budget", "authority", "need", "timeline")]):
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "BANT（商談情報）"
            t = slide.shapes.title.text_frame.paragraphs[0]
            if t.runs:
                set_font(t.runs[0], 30, True, PRIMARY_RGB)
            t.alignment = PP_ALIGN.LEFT
        box = ensure_textbox(slide, 0.6, 1.9, 9.0, 5.0)
        tf = box.text_frame
        tf.clear()
        rows = [
            ("Budget（予算）", bant.get("budget", "")),
            ("Authority（決裁）", bant.get("authority", "")),
            ("Need（ニーズ）", bant.get("need", "")),
            ("Timeline（時期）", bant.get("timeline", "")),
        ]
        for i, (k, v) in enumerate(rows):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = k
            if p.runs: set_font(p.runs[0], 18, True, PRIMARY_RGB)
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.1
            p2 = tf.add_paragraph()
            p2.text = f"{v or '—'}"
            if p2.runs: set_font(p2.runs[0], 18, False, TEXT_RGB)
            p2.space_after = Pt(8)

    def add_next_actions_slide(items: list):
        if not items:
            return
        # ネクストアクションはワンスライド・ワンメッセージに沿って可能なら個別スライド化
        for it in items:
            slide = prs.slides.add_slide(get_layout(1))
            if slide.shapes.title:
                slide.shapes.title.text = "ネクストアクション"
                t = slide.shapes.title.text_frame.paragraphs[0]
                if t.runs:
                    set_font(t.runs[0], 30, True, PRIMARY_RGB)
                t.alignment = PP_ALIGN.LEFT
            promote_key_message(slide, shorten_bullet(it, max_chars=100))

    def add_summary_slide(items: list):
        if not items:
            return
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "まとめ"
            t = slide.shapes.title.text_frame.paragraphs[0]
            if t.runs:
                set_font(t.runs[0], 30, True, PRIMARY_RGB)
            t.alignment = PP_ALIGN.LEFT
        add_bullets_block(ensure_textbox(slide, 0.6, 1.9, 9.0, 4.8), [shorten_bullet(i) for i in items], 20, "●")

    # ===== 実装：スライド生成フロー =====
    title = parsed.get("title") or "打ち合わせ要約"
    company = parsed.get("company_name") or ""
    date = parsed.get("meeting_date") or ""
    subtitle = " / ".join([s for s in [company, f"{date} 実施" if date else ""] if s])

    add_title_slide(title, subtitle)
    add_agenda_slide(parsed.get("agenda", []))

    # セクションごと
    sections = parsed.get("sections", [])
    for s in sections:
        stitle = s.get("title") or ""
        bullets = [shorten_bullet(b) for b in (s.get("bullets") or [])]
        notes = s.get("notes") or []
        # セクション見出しスライド（大見出しに見える場合のみ）
        if stitle and (len(bullets) == 0 and len(notes) == 0):
            add_section_divider(stitle)
        # 通常コンテンツ
        if bullets:
            add_content_slide(stitle or "セクション", bullets, notes)
        elif notes:
            add_content_slide(stitle or "セクション", [], notes)

    # BANT／課題・ニーズ／ネクストアクション／まとめ
    add_bant_slide(parsed.get("bant", {}))

    challenges = [shorten_bullet(x) for x in parsed.get("challenges", [])]
    needs = [shorten_bullet(x) for x in parsed.get("needs", [])]
    if challenges or needs:
        slide = prs.slides.add_slide(get_layout(1))
        if slide.shapes.title:
            slide.shapes.title.text = "課題とニーズ"
            t = slide.shapes.title.text_frame.paragraphs[0]
            if t.runs:
                set_font(t.runs[0], 30, True, PRIMARY_RGB)
            t.alignment = PP_ALIGN.LEFT
        add_bullets_block(ensure_textbox(slide, 0.6, 1.9, 4.4, 4.8), [shorten_bullet(i) for i in challenges], 20)
        add_bullets_block(ensure_textbox(slide, 5.2, 1.9, 4.4, 4.8), [shorten_bullet(i) for i in needs], 20)

    add_next_actions_slide(parsed.get("next_actions", []))
    add_summary_slide(parsed.get("summary", []))

    # 末尾に Thanks スライド
    slide = prs.slides.add_slide(get_layout(1))
    if slide.shapes.title:
        slide.shapes.title.text = "ご清聴ありがとうございました"
        t = slide.shapes.title.text_frame.paragraphs[0]
        if t.runs:
            set_font(t.runs[0], 32, True, PRIMARY_RGB)
        t.alignment = PP_ALIGN.LEFT

    # 出力
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ===== Flask ルーティング =====
@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    try:
        minutes_text = request.form.get("minutes_text", "")
        
        if not minutes_text.strip():
            return render_template("index.html", error="議事録テキストを入力してください。")
        
        # 議事録を解析
        parsed = parse_meeting_minutes(minutes_text)
        
        # PowerPointファイルを生成
        ppt_file = create_meeting_summary_ppt(parsed)

        # ファイル名を生成
        company_name = parsed.get("company_name", "meeting") or "meeting"
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        fname = f"{safe_ascii_filename(company_name)}_summary_{timestamp}.pptx"
        
        # 一時ファイルとして保存
        file_path = os.path.join(TEMP_DIR, fname)
        with open(file_path, 'wb') as f:
            f.write(ppt_file.getvalue())
        
        # セッションに情報を保存
        session['last_generated'] = {
            'filename': fname,
            'company_name': company_name,
            'slide_count': len(Presentation(file_path).slides),
            'analysis_summary': f"{company_name}様の議事録を解析し、{len(parsed.get('sections', []))}セクションのスライドを生成"
        }
        
        # 成功画面にリダイレクト
        return redirect(url_for('success'))
        
    except ValueError as e:
        return render_template("index.html", error=str(e))
    except Exception as e:
        print(f"[ERROR] Generation failed: {e}")
        return render_template("index.html", error="PowerPoint生成中にエラーが発生しました。Azure OpenAIの設定を確認してください。")

@app.route("/success")
def success():
    if 'last_generated' not in session:
        return redirect(url_for('index'))
    
    info = session['last_generated']
    return render_template("success.html", 
                         filename=info['filename'],
                         slide_count=info['slide_count'],
                         analysis_summary=info['analysis_summary'])

@app.route("/download/<filename>")
def download_file(filename):
    try:
        file_path = os.path.join(TEMP_DIR, filename)
        if not os.path.exists(file_path):
            return render_template("index.html", error="ファイルが見つかりません。")
        
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        print(f"[ERROR] Download failed: {e}")
        return render_template("index.html", error="ファイルのダウンロードに失敗しました。")

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5001))
    app.run(debug=False, host="0.0.0.0", port=port)
